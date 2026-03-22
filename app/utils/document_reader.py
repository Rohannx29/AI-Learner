import pdfplumber
import pytesseract
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
from PIL import Image
import cv2
import numpy as np

from app.config import settings
from app.utils.math_ocr import extract_math_from_image

if settings.TESSERACT_PATH:
    pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_PATH


def preprocess_image(image):
    img = np.array(image)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    blur = cv2.GaussianBlur(gray, (5, 5), 0)
    thresh = cv2.adaptiveThreshold(
        blur, 255,
        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
        cv2.THRESH_BINARY,
        11, 2
    )
    return thresh


def read_image(file_path: str) -> str:
    image = Image.open(file_path)
    processed = preprocess_image(image)
    text = pytesseract.image_to_string(processed)
    math = extract_math_from_image(file_path)
    return text + "\n" + math


def read_pdf(file_path: str) -> str:
    text = ""

    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"

    if not text.strip():
        kwargs = {}
        if settings.POPPLER_PATH:
            kwargs["poppler_path"] = settings.POPPLER_PATH

        images = convert_from_path(file_path, **kwargs)

        for img in images:
            processed = preprocess_image(img)
            ocr_text = pytesseract.image_to_string(processed)
            math = extract_math_from_image(img)
            text += ocr_text + "\n" + math

    return text


def read_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join(para.text for para in doc.paragraphs)


def read_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def read_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()